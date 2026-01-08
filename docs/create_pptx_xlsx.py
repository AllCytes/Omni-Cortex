"""
OmniCortex PowerPoint & Excel Generator
Creates: OmniCortex_Overview.pptx and OmniCortex_ToolReference.xlsx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

# Color palette (light theme)
PRIMARY = "4A90D9"      # Blue
SECONDARY = "7C3AED"    # Purple
ACCENT = "10B981"       # Green
DARK_TEXT = "1F2937"
LIGHT_BG = "F8FAFC"
WHITE = "FFFFFF"


def hex_to_rgb(hex_color):
    """Convert hex color to RGBColor"""
    return RGBColor(int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def create_powerpoint():
    """Create OmniCortex Overview PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    # === SLIDE 1: Title ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = hex_to_rgb(PRIMARY)
    bg.line.fill.background()

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "OmniCortex"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(WHITE)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "Universal Memory System for Claude Code"
    p.font.size = Pt(28)
    p.font.color.rgb = hex_to_rgb(WHITE)
    p.alignment = PP_ALIGN.CENTER

    # Version
    ver = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.333), Inches(0.5))
    tf = ver.text_frame
    p = tf.paragraphs[0]
    p.text = "v1.0.5 | pip install omni-cortex"
    p.font.size = Pt(18)
    p.font.color.rgb = hex_to_rgb(WHITE)
    p.alignment = PP_ALIGN.CENTER

    # === SLIDE 2: What is OmniCortex ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = hex_to_rgb(PRIMARY)
    header.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "What is OmniCortex?"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(WHITE)

    # Content
    content = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(5.5))
    tf = content.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "A persistent memory system that stores what you learn, tracks what you do, and helps you never repeat the same mistakes twice."
    p.font.size = Pt(24)
    p.font.color.rgb = hex_to_rgb(DARK_TEXT)
    p.space_after = Pt(20)

    features = [
        ("18 Tools", "Memory, activity, session, and global search tools"),
        ("Dual Storage", "Activity log + Knowledge store"),
        ("Smart Search", "FTS5 keyword + Semantic + Hybrid modes"),
        ("Session Continuity", '"Last time you were working on..."'),
        ("Cross-Project", "Global index searches all your projects"),
        ("Web Dashboard", "Full GUI for browsing and analytics"),
    ]

    for title_text, desc in features:
        p = tf.add_paragraph()
        p.text = f"{title_text}: {desc}"
        p.font.size = Pt(20)
        p.font.color.rgb = hex_to_rgb(DARK_TEXT)
        p.space_before = Pt(12)
        p.level = 0

    # === SLIDE 3: Core Features ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = hex_to_rgb(SECONDARY)
    header.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Core Features"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(WHITE)

    # Two columns
    left_features = [
        "Memory Tools (6)",
        "  - cortex_remember",
        "  - cortex_recall",
        "  - cortex_list_memories",
        "  - cortex_update_memory",
        "  - cortex_forget",
        "  - cortex_link_memories",
        "",
        "Activity Tools (3)",
        "  - cortex_log_activity",
        "  - cortex_get_activities",
        "  - cortex_get_timeline",
    ]

    right_features = [
        "Session Tools (3)",
        "  - cortex_start_session",
        "  - cortex_end_session",
        "  - cortex_get_session_context",
        "",
        "Global Tools (3)",
        "  - cortex_global_search",
        "  - cortex_global_stats",
        "  - cortex_sync_to_global",
        "",
        "Utility Tools (3)",
        "  - cortex_list_tags",
        "  - cortex_review_memories",
        "  - cortex_export",
    ]

    left_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(5.5), Inches(5.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(left_features):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16) if line.startswith("  ") else Pt(20)
        p.font.bold = not line.startswith("  ") and line != ""
        p.font.color.rgb = hex_to_rgb(SECONDARY if not line.startswith("  ") and line else DARK_TEXT)

    right_box = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.5), Inches(5.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(right_features):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16) if line.startswith("  ") else Pt(20)
        p.font.bold = not line.startswith("  ") and line != ""
        p.font.color.rgb = hex_to_rgb(SECONDARY if not line.startswith("  ") and line else DARK_TEXT)

    # === SLIDE 4: Web Dashboard ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = hex_to_rgb(ACCENT)
    header.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Web Dashboard"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(WHITE)

    # Command box
    cmd_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(1.5), Inches(7), Inches(0.8))
    cmd_box.fill.solid()
    cmd_box.fill.fore_color.rgb = hex_to_rgb(DARK_TEXT)
    cmd_text = slide.shapes.add_textbox(Inches(3), Inches(1.6), Inches(7), Inches(0.6))
    tf = cmd_text.text_frame
    p = tf.paragraphs[0]
    p.text = "omni-cortex-dashboard"
    p.font.size = Pt(24)
    p.font.name = "Courier New"
    p.font.color.rgb = hex_to_rgb(ACCENT)
    p.alignment = PP_ALIGN.CENTER

    # 6 Tabs
    tabs = [
        ("Memories", "Browse, search, filter, edit memories"),
        ("Activity", "Complete audit trail of tool usage"),
        ("Statistics", "Heatmaps, charts, distributions"),
        ("Review", "Mark stale memories fresh/outdated"),
        ("Graph", "D3.js relationship visualization"),
        ("Ask AI", "Gemini-powered natural language queries"),
    ]

    y_pos = 2.6
    for name, desc in tabs:
        tab_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(y_pos), Inches(11.333), Inches(0.7))
        tab_shape.fill.solid()
        tab_shape.fill.fore_color.rgb = hex_to_rgb("E8F5E9")
        tab_shape.line.color.rgb = hex_to_rgb(ACCENT)

        tab_text = slide.shapes.add_textbox(Inches(1.2), Inches(y_pos + 0.15), Inches(11), Inches(0.5))
        tf = tab_text.text_frame
        p = tf.paragraphs[0]
        p.text = f"{name} - {desc}"
        p.font.size = Pt(18)
        p.font.color.rgb = hex_to_rgb(DARK_TEXT)

        y_pos += 0.75

    # === SLIDE 5: Installation ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = hex_to_rgb(PRIMARY)
    header.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Quick Start"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(WHITE)

    steps = [
        ("1. Install", "pip install omni-cortex"),
        ("2. Setup", "omni-cortex-setup"),
        ("3. Restart", "Restart Claude Code"),
        ("4. Dashboard (optional)", "omni-cortex-dashboard"),
    ]

    y_pos = 1.8
    for step_title, cmd in steps:
        # Step title
        step_text = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(4), Inches(0.5))
        tf = step_text.text_frame
        p = tf.paragraphs[0]
        p.text = step_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(PRIMARY)

        # Command box
        cmd_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(y_pos - 0.1), Inches(7), Inches(0.7))
        cmd_shape.fill.solid()
        cmd_shape.fill.fore_color.rgb = hex_to_rgb(DARK_TEXT)

        cmd_text = slide.shapes.add_textbox(Inches(5.2), Inches(y_pos), Inches(6.6), Inches(0.5))
        tf = cmd_text.text_frame
        p = tf.paragraphs[0]
        p.text = cmd
        p.font.size = Pt(18)
        p.font.name = "Courier New"
        p.font.color.rgb = hex_to_rgb(ACCENT)

        y_pos += 1.2

    # Semantic search note
    note = slide.shapes.add_textbox(Inches(1), Inches(6.3), Inches(11), Inches(0.5))
    tf = note.text_frame
    p = tf.paragraphs[0]
    p.text = "For AI-powered semantic search: pip install omni-cortex[semantic]"
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = hex_to_rgb(SECONDARY)

    # === SLIDE 6: Thank You ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = hex_to_rgb(SECONDARY)
    bg.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Get Started Today"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(WHITE)
    p.alignment = PP_ALIGN.CENTER

    links = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(2))
    tf = links.text_frame
    p = tf.paragraphs[0]
    p.text = "github.com/AllCytes/Omni-Cortex"
    p.font.size = Pt(24)
    p.font.color.rgb = hex_to_rgb(WHITE)
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "pypi.org/project/omni-cortex"
    p.font.size = Pt(24)
    p.font.color.rgb = hex_to_rgb(WHITE)
    p.alignment = PP_ALIGN.CENTER

    prs.save('D:/Projects/omni-cortex/docs/OmniCortex_Overview.pptx')
    print("[OK] Created: OmniCortex_Overview.pptx")


def create_excel():
    """Create OmniCortex Tool Reference Excel spreadsheet"""
    wb = Workbook()
    ws = wb.active
    ws.title = "Tool Reference"
    ws.sheet_properties.tabColor = PRIMARY

    # Styles
    thin = Side(style='thin', color='CCCCCC')
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    header_fill = PatternFill(start_color=PRIMARY, end_color=PRIMARY, fill_type='solid')
    header_font = Font(bold=True, color=WHITE, size=11)

    alt_fill = PatternFill(start_color='F3F4F6', end_color='F3F4F6', fill_type='solid')

    # Title
    ws.merge_cells('A1:F1')
    ws['A1'] = 'OmniCortex Tool Reference'
    ws['A1'].font = Font(bold=True, size=18, color=PRIMARY)
    ws['A1'].alignment = Alignment(horizontal='center', vertical='center')
    ws.row_dimensions[1].height = 30

    ws.merge_cells('A2:F2')
    ws['A2'] = 'v1.0.5 | 18 Tools | pip install omni-cortex'
    ws['A2'].font = Font(size=11, color='666666')
    ws['A2'].alignment = Alignment(horizontal='center')

    # Headers
    headers = ['Tool Name', 'Category', 'Description', 'Key Parameters', 'Return Type', 'Notes']
    for col, header in enumerate(headers, 1):
        cell = ws.cell(row=4, column=col, value=header)
        cell.fill = header_fill
        cell.font = header_font
        cell.border = border
        cell.alignment = Alignment(horizontal='center', vertical='center')
    ws.row_dimensions[4].height = 25

    # Tool data
    tools = [
        ('cortex_remember', 'Memory', 'Store information with auto-categorization', 'content, context, tags, type, importance', 'memory_id', 'Auto-detects type if not specified'),
        ('cortex_recall', 'Memory', 'Search memories by keyword or meaning', 'query, search_mode, type_filter, limit', 'memories[]', 'Modes: keyword, semantic, hybrid'),
        ('cortex_list_memories', 'Memory', 'List all memories with filters', 'type_filter, sort_by, limit, offset', 'memories[]', 'Supports pagination'),
        ('cortex_update_memory', 'Memory', 'Update an existing memory', 'id, content, add_tags, status', 'updated memory', 'Status: fresh, outdated, archived'),
        ('cortex_forget', 'Memory', 'Permanently delete a memory', 'id, confirm=true', 'confirmation', 'Requires confirm=true'),
        ('cortex_link_memories', 'Memory', 'Create relationship between memories', 'source_id, target_id, relationship_type', 'link_id', 'Types: related_to, supersedes, etc.'),
        ('cortex_log_activity', 'Activity', 'Manually log an activity', 'event_type, tool_name, success', 'activity_id', 'Usually auto-logged via hooks'),
        ('cortex_get_activities', 'Activity', 'Query the activity log', 'session_id, tool_name, since, limit', 'activities[]', 'Filter by session or tool'),
        ('cortex_get_timeline', 'Activity', 'Chronological view of activities', 'hours, include_activities, group_by', 'timeline', 'Group by: hour, day, session'),
        ('cortex_start_session', 'Session', 'Start a new work session', 'provide_context, context_depth', 'session_id + context', 'Returns previous session context'),
        ('cortex_end_session', 'Session', 'End session with summary', 'session_id, summary, key_learnings', 'session summary', 'Summary auto-generated if omitted'),
        ('cortex_get_session_context', 'Session', 'Get context from past sessions', 'session_count, include_decisions', 'context summary', 'For resuming work'),
        ('cortex_list_tags', 'Utility', 'List all tags with counts', 'min_count, limit', 'tags[]', 'Shows tag usage frequency'),
        ('cortex_review_memories', 'Utility', 'Review memory freshness', 'action, days_threshold, memory_ids', 'review results', 'Actions: list, mark_fresh, archive'),
        ('cortex_export', 'Utility', 'Export memories and activities', 'format, output_path, since', 'export data', 'Formats: markdown, json, sqlite'),
        ('cortex_global_search', 'Global', 'Search across all projects', 'query, project_filter, limit', 'memories[]', 'Searches ~/.omni-cortex/global.db'),
        ('cortex_global_stats', 'Global', 'Get global index statistics', '(none)', 'statistics', 'Memory counts by project/type'),
        ('cortex_sync_to_global', 'Global', 'Sync project to global index', 'full_sync', 'sync count', 'Usually automatic'),
    ]

    # Category colors
    cat_colors = {
        'Memory': 'E3F2FD',
        'Activity': 'F3E5F5',
        'Session': 'E8F5E9',
        'Utility': 'FFF3E0',
        'Global': 'E0F7FA',
    }

    for row_idx, (name, cat, desc, params, ret, notes) in enumerate(tools, 5):
        row_data = [name, cat, desc, params, ret, notes]
        for col_idx, value in enumerate(row_data, 1):
            cell = ws.cell(row=row_idx, column=col_idx, value=value)
            cell.border = border
            cell.font = Font(size=10)

            if col_idx == 1:  # Tool name
                cell.font = Font(size=10, name='Courier New', color=PRIMARY)
            elif col_idx == 2:  # Category
                cell.fill = PatternFill(start_color=cat_colors.get(cat, 'FFFFFF'),
                                       end_color=cat_colors.get(cat, 'FFFFFF'), fill_type='solid')
                cell.alignment = Alignment(horizontal='center')

            if row_idx % 2 == 0 and col_idx not in [1, 2]:
                cell.fill = alt_fill

    # Column widths
    widths = {'A': 25, 'B': 12, 'C': 40, 'D': 45, 'E': 18, 'F': 35}
    for col, width in widths.items():
        ws.column_dimensions[col].width = width

    # Freeze header
    ws.freeze_panes = 'A5'

    # Add auto-filter
    ws.auto_filter.ref = f"A4:F{4 + len(tools)}"

    # === SECOND SHEET: Memory Types ===
    ws2 = wb.create_sheet("Memory Types")
    ws2.sheet_properties.tabColor = SECONDARY

    ws2.merge_cells('A1:C1')
    ws2['A1'] = 'Memory Type Detection'
    ws2['A1'].font = Font(bold=True, size=16, color=SECONDARY)

    type_headers = ['Type', 'Auto-detected When...', 'Example']
    for col, header in enumerate(type_headers, 1):
        cell = ws2.cell(row=3, column=col, value=header)
        cell.fill = PatternFill(start_color=SECONDARY, end_color=SECONDARY, fill_type='solid')
        cell.font = header_font
        cell.border = border

    types = [
        ('solution', 'Contains "fix", "resolved", "solution"', 'Fixed the timeout by increasing buffer size'),
        ('warning', 'Contains "warning", "avoid", "never"', 'Never use deprecated API v1'),
        ('config', 'Contains "config", "setting", "environment"', 'Set DEBUG=true in .env'),
        ('troubleshooting', 'Contains "debug", "troubleshoot"', 'Debug by checking network tab'),
        ('code', 'Contains "function", "class", "algorithm"', 'Use async/await pattern for API calls'),
        ('error', 'Contains "error", "exception", "failed"', 'ValueError when input is None'),
        ('command', 'Contains "run", "execute", "command"', 'Run npm install --legacy-peer-deps'),
        ('concept', 'Contains "what is", "definition"', 'MCP is Model Context Protocol'),
        ('decision', 'Contains "decided", "chose", "architecture"', 'Chose SQLite over PostgreSQL'),
        ('tip', 'Contains "tip", "trick", "best practice"', 'Tip: Use batch updates for performance'),
        ('general', 'Default type', 'Any other information'),
    ]

    for row_idx, (typ, detection, example) in enumerate(types, 4):
        ws2.cell(row=row_idx, column=1, value=typ).font = Font(bold=True, color=SECONDARY)
        ws2.cell(row=row_idx, column=2, value=detection)
        ws2.cell(row=row_idx, column=3, value=example).font = Font(italic=True, color='666666')
        for col in range(1, 4):
            ws2.cell(row=row_idx, column=col).border = border
            if row_idx % 2 == 0:
                ws2.cell(row=row_idx, column=col).fill = alt_fill

    ws2.column_dimensions['A'].width = 18
    ws2.column_dimensions['B'].width = 45
    ws2.column_dimensions['C'].width = 45

    wb.save('D:/Projects/omni-cortex/docs/OmniCortex_ToolReference.xlsx')
    print("[OK] Created: OmniCortex_ToolReference.xlsx")


if __name__ == "__main__":
    print("Creating OmniCortex PowerPoint and Excel files...")
    create_powerpoint()
    create_excel()
    print("\nAll files created successfully!")
